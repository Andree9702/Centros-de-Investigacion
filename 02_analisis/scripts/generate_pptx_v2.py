"""
generate_pptx_v2.py
Presentacion_Centros_UTMACH_Final.pptx — 10 slides de alto impacto.
Reglas: max 4 bullets/slide, titulos >=36pt, cuerpo >=20pt, imagenes 6 in centradas.
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── RUTAS ─────────────────────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS_DIR = os.path.join(BASE_DIR, "resultados", "figuras")
OUTPUT   = os.path.join(BASE_DIR, "resultados", "Presentacion_Centros_UTMACH_Final.pptx")

F = {
    "clasif":    os.path.join(FIGS_DIR, "01_clasificacion_produccion.png"),
    "impacto":   os.path.join(FIGS_DIR, "02_impacto_por_centro.png"),
    "cuartiles": os.path.join(FIGS_DIR, "03_cuartiles_por_centro.png"),
    "lineas":    os.path.join(FIGS_DIR, "04_top_lineas_por_tipo.png"),
    "ippc":      os.path.join(FIGS_DIR, "05_masa_critica_ippc.png"),
    "dashboard": os.path.join(FIGS_DIR, "06_dashboard_ejecutivo.png"),
}

# ── PALETA ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0B, 0x18, 0x2A)   # azul medianoche
PANEL   = RGBColor(0x12, 0x2B, 0x48)   # panel más claro
ACCENT  = RGBColor(0x00, 0x8C, 0xFF)   # azul eléctrico
GREEN   = RGBColor(0x00, 0xC9, 0x8E)   # verde esmeralda
GOLD    = RGBColor(0xF5, 0xA6, 0x23)   # dorado
RED     = RGBColor(0xE0, 0x3B, 0x3B)   # rojo alerta
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xC8, 0xD8, 0xE8)   # gris claro texto

W = Inches(13.33)
H = Inches(7.5)

# ── HELPERS MÍNIMOS ───────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    sld = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sld.background.fill.solid()
    sld.background.fill.fore_color.rgb = BG
    return sld


def rect(sld, x, y, w, h, color, line=False):
    s = sld.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s


def txt(sld, text, x, y, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = sld.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


def bullet_box(sld, items, x, y, w, h, size=20, color=LGRAY,
               spacing=Pt(10), marker="  "):
    """
    items: list of str  OR  list of (str, RGBColor)
    """
    box = sld.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, clr = item
        else:
            text, clr = item, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = spacing
        r = p.add_run()
        r.text           = marker + text
        r.font.size      = Pt(size)
        r.font.color.rgb = clr
    return box


def insert_img(sld, path, img_w=Inches(6.0), top=Inches(1.3)):
    """Inserta imagen centrada horizontalmente con ancho fijo."""
    try:
        from PIL import Image as PILImage
        with PILImage.open(path) as im:
            iw, ih = im.size
        h = img_w * (ih / iw)
    except Exception:
        h = img_w * 0.6
    left = (W - img_w) / 2
    sld.shapes.add_picture(path, left, top, width=img_w, height=h)
    return top + h   # retorna y-bottom de la imagen


def title_bar(sld, title, sub=None, bar_h=Inches(1.05)):
    """Barra de título con línea de acento superior."""
    rect(sld, 0, 0, W, bar_h, PANEL)
    rect(sld, 0, 0, W, Inches(0.055), ACCENT)
    txt(sld, title,
        Inches(0.45), Inches(0.1), W - Inches(0.9), Inches(0.65),
        size=36, bold=True, color=WHITE)
    if sub:
        txt(sld, sub,
            Inches(0.45), Inches(0.72), W - Inches(0.9), Inches(0.32),
            size=16, color=ACCENT)
    return bar_h


def footer(sld, page, total=10,
           text="UTMACH  |  DIDI  |  Marzo 2026  |  2.919 articulos"):
    rect(sld, 0, H - Inches(0.28), W, Inches(0.28), PANEL)
    txt(sld, text,
        Inches(0.4), H - Inches(0.26), W * 0.7, Inches(0.24),
        size=9, color=LGRAY)
    txt(sld, f"{page}/{total}",
        W - Inches(0.9), H - Inches(0.26), Inches(0.8), Inches(0.24),
        size=9, color=LGRAY, align=PP_ALIGN.RIGHT)


# ── SLIDES ────────────────────────────────────────────────────────────────────

# ── S1: TÍTULO ────────────────────────────────────────────────────────────────
def s01_titulo(prs):
    sld = add_slide(prs)
    rect(sld, 0, 0, Inches(0.3), H, ACCENT)

    txt(sld, "De la Intuicion a la Evidencia",
        Inches(0.55), Inches(1.2), Inches(11.5), Inches(1.2),
        size=46, bold=True, color=WHITE)

    txt(sld, "Data-Driven",
        Inches(0.55), Inches(2.3), Inches(11.5), Inches(0.9),
        size=46, bold=True, color=ACCENT)

    rect(sld, Inches(0.55), Inches(3.35), Inches(9.0), Inches(0.05), GOLD)

    txt(sld, "Arquitectura de Centros de Investigacion  |  UTMACH",
        Inches(0.55), Inches(3.55), Inches(11.5), Inches(0.55),
        size=22, color=LGRAY)

    txt(sld, "Corpus: 2.919 articulos  |  NLP Semantico + K-Means  |  Sistema de Dos Vias",
        Inches(0.55), Inches(4.15), Inches(11.5), Inches(0.45),
        size=18, color=LGRAY, italic=True)

    # Equipo
    rect(sld, Inches(0.55), Inches(5.2), Inches(10.5), Inches(1.4), PANEL)
    txt(sld, "PhD. Ivan Ramirez — Director DIDI   |   "
             "MSc. Andree Vitonera   |   MSc. Luiggi Solano",
        Inches(0.75), Inches(5.55), Inches(10.1), Inches(0.5),
        size=16, color=LGRAY, align=PP_ALIGN.CENTER)

    footer(sld, 1)
    return sld


# ── S2: EL PROBLEMA ───────────────────────────────────────────────────────────
def s02_problema(prs):
    sld = add_slide(prs)
    title_bar(sld, "El Problema",
              "Por que los metodos tradicionales producen centros ineficientes")

    # Dos columnas: Errores | Costo
    col_y  = Inches(1.2)
    col_h  = Inches(5.0)
    col_w  = Inches(5.9)

    # Col izquierda
    rect(sld, Inches(0.4), col_y, col_w, Inches(0.5), RED)
    txt(sld, "ERRORES COMUNES",
        Inches(0.55), col_y + Inches(0.07), col_w, Inches(0.4),
        size=20, bold=True, color=WHITE)

    errores = [
        "Centros por cuotas de facultad = sesgo politico",
        "Centros por palabras repetidas = falsa precision",
        "Experimental + Social en el mismo modelo = desperdicio presupuestario",
        "Prioridades a priori = analisis viciado",
    ]
    y = col_y + Inches(0.65)
    for e in errores:
        rect(sld, Inches(0.4), y, col_w, Inches(0.82), PANEL)
        txt(sld, e, Inches(0.55), y + Inches(0.13),
            col_w - Inches(0.25), Inches(0.62), size=20, color=LGRAY)
        y += Inches(0.9)

    # Col derecha
    x2 = Inches(7.0)
    rect(sld, x2, col_y, col_w, Inches(0.5), GOLD)
    txt(sld, "EL COSTO REAL",
        x2 + Inches(0.15), col_y + Inches(0.07), col_w, Inches(0.4),
        size=20, bold=True, color=BG)

    costos = [
        "Laboratorios financiando proyectos de escritorio",
        "Observatorios pagando reactivos que no usan",
        "Investigadores elite fuera de su area de impacto",
        "Recursos diluidos  —  impacto institucional minimo",
    ]
    y = col_y + Inches(0.65)
    for c in costos:
        rect(sld, x2, y, col_w, Inches(0.82), PANEL)
        txt(sld, c, x2 + Inches(0.15), y + Inches(0.13),
            col_w - Inches(0.25), Inches(0.62), size=20, color=LGRAY)
        y += Inches(0.9)

    footer(sld, 2)
    return sld


# ── S3: DESCUBRIMIENTO ADN ────────────────────────────────────────────────────
def s03_adn(prs):
    sld = add_slide(prs)
    title_bar(sld, "El Gran Descubrimiento: El ADN Cientifico de la UTMACH")

    img_bottom = insert_img(sld, F["clasif"], img_w=Inches(6.5), top=Inches(1.15))

    # Panel de texto a la derecha de la imagen
    px = (W + Inches(6.5)) / 2 + Inches(0.3)
    pw = W - px - Inches(0.3)
    py = Inches(1.2)

    kpis = [
        ("47%",  "Experimental",         ACCENT),
        ("41%",  "Social / Humanistica", GREEN),
        ("11.6%","Revisiones Bibl.",     GOLD),
    ]
    y = py
    for val, label, clr in kpis:
        rect(sld, px, y, pw, Inches(1.45), PANEL)
        rect(sld, px, y, Inches(0.08), Inches(1.45), clr)
        txt(sld, val, px + Inches(0.18), y + Inches(0.08),
            pw - Inches(0.25), Inches(0.72), size=36, bold=True, color=clr)
        txt(sld, label, px + Inches(0.18), y + Inches(0.82),
            pw - Inches(0.25), Inches(0.5), size=18, color=LGRAY)
        y += Inches(1.6)

    # Mensaje clave
    rect(sld, Inches(0.4), H - Inches(1.1), W - Inches(0.8), Inches(0.8), PANEL)
    rect(sld, Inches(0.4), H - Inches(1.1), Inches(0.08), Inches(0.8), GOLD)
    txt(sld,
        "47% Experimental vs 41% Social  ->  La UTMACH exige DOS modelos de gestion financiera distintos.",
        Inches(0.6), H - Inches(1.0), W - Inches(1.2), Inches(0.65),
        size=20, bold=True, color=GOLD)

    footer(sld, 3)
    return sld


# ── S4: SISTEMA DE DOS VIAS ───────────────────────────────────────────────────
def s04_dos_vias(prs):
    sld = add_slide(prs)
    title_bar(sld, "El Sistema de Dos Vias",
              "Una arquitectura financiera diferenciada — no un organigrama")

    col_w = Inches(5.85)
    cy    = Inches(1.2)

    # VIA A
    rect(sld, Inches(0.4), cy, col_w, Inches(0.55), ACCENT)
    txt(sld, "VIA A — CENTROS CIENTIFICOS EXPERIMENTALES",
        Inches(0.55), cy + Inches(0.07), col_w, Inches(0.44),
        size=18, bold=True, color=WHITE)

    rows_a = [
        ("Base",          "1.371 articulos  (47%)"),
        ("Infraestructura","Laboratorios, reactivos, equipos especializados"),
        ("Presupuesto",   "ALTO — inversion de capital a largo plazo"),
        ("Estructuras",   "5 Centros de Investigacion"),
    ]
    y = cy + Inches(0.6)
    for campo, valor in rows_a:
        rect(sld, Inches(0.4), y, col_w, Inches(0.78), PANEL)
        txt(sld, campo + ":", Inches(0.55), y + Inches(0.1),
            Inches(2.0), Inches(0.55), size=20, bold=True, color=ACCENT)
        txt(sld, valor, Inches(2.55), y + Inches(0.1),
            col_w - Inches(2.25), Inches(0.55), size=20, color=LGRAY)
        y += Inches(0.85)

    # VIA B
    x2 = Inches(7.1)
    rect(sld, x2, cy, col_w, Inches(0.55), GREEN)
    txt(sld, "VIA B — OBSERVATORIOS SOCIALES Y HUMANISTICOS",
        x2 + Inches(0.15), cy + Inches(0.07), col_w, Inches(0.44),
        size=18, bold=True, color=BG)

    rows_b = [
        ("Base",          "1.209 articulos  (41.4%)"),
        ("Infraestructura","Software, bases de datos, plataformas digitales"),
        ("Presupuesto",   "MODERADO — gestion agil y flexible"),
        ("Estructuras",   "5 Observatorios Academicos"),
    ]
    y = cy + Inches(0.6)
    for campo, valor in rows_b:
        rect(sld, x2, y, col_w, Inches(0.78), PANEL)
        txt(sld, campo + ":", x2 + Inches(0.15), y + Inches(0.1),
            Inches(2.0), Inches(0.55), size=20, bold=True, color=GREEN)
        txt(sld, valor, x2 + Inches(2.15), y + Inches(0.1),
            col_w - Inches(2.25), Inches(0.55), size=20, color=LGRAY)
        y += Inches(0.85)

    # Divisor
    rect(sld, Inches(6.6), cy, Inches(0.06), Inches(4.0), GOLD)

    footer(sld, 4)
    return sld


# ── S5: VIA A CENTROS ─────────────────────────────────────────────────────────
def s05_via_a(prs):
    sld = add_slide(prs)
    title_bar(sld, "VIA A — Los 5 Centros Cientificos Experimentales",
              "Impacto ponderado: la metrica que supera al conteo de articulos")

    img_bottom = insert_img(sld, F["impacto"], img_w=Inches(6.2), top=Inches(1.12))

    # Tabla de centros a la derecha
    centros = [
        ("Ingenieria y Analisis Cuantitativo (II)", "366", "204.5", GOLD),
        ("Ciencias Quimicas y Ambientales",         "278", "138.2", ACCENT),
        ("Ingenieria y Analisis Cuantitativo",      "284", "137.2", ACCENT),
        ("Salud Integral y Biociencias",            "265", "126.6", GREEN),
        ("Agroalimentaria y Sostenibilidad",        "178", " 91.2", GREEN),
    ]
    px = (W + Inches(6.2)) / 2 + Inches(0.2)
    pw = W - px - Inches(0.25)
    y  = Inches(1.15)

    # Header
    rect(sld, px, y, pw, Inches(0.42), ACCENT)
    for label, xoff in [("Arts.", Inches(0.0)), ("Impacto", Inches(1.45))]:
        txt(sld, label, px + xoff + pw * 0.55, y + Inches(0.05),
            Inches(1.3), Inches(0.34), size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
    y += Inches(0.42)

    for nombre, arts, imp, clr in centros:
        rect(sld, px, y, pw, Inches(0.96), PANEL)
        rect(sld, px, y, Inches(0.06), Inches(0.96), clr)
        txt(sld, nombre, px + Inches(0.12), y + Inches(0.06),
            pw * 0.52, Inches(0.84), size=13, color=WHITE)
        txt(sld, arts, px + pw * 0.55, y + Inches(0.2),
            Inches(1.2), Inches(0.55), size=18, bold=True, color=LGRAY,
            align=PP_ALIGN.CENTER)
        txt(sld, imp,  px + pw * 0.55 + Inches(1.3), y + Inches(0.2),
            Inches(1.3), Inches(0.55), size=18, bold=True, color=clr,
            align=PP_ALIGN.CENTER)
        y += Inches(1.01)

    footer(sld, 5)
    return sld


# ── S6: VIA B OBSERVATORIOS ───────────────────────────────────────────────────
def s06_via_b(prs):
    sld = add_slide(prs)
    title_bar(sld, "VIA B — Los 5 Observatorios Sociales y Humanisticos",
              "Gestion agil, alto volumen, transferencia al sector publico")

    obs = [
        ("Obs. de Economia, Empresa e Innovacion",         "547", "269.6", "0.493",
         "El mas grande de la institucion — 18.7% de toda la produccion UTMACH", GOLD),
        ("Obs. Multidisciplinario de Ciencias y Sociedad", "200", "128.3", "0.642",
         "Mayor promedio de impacto de la Via B  (41 articulos Q1)", GREEN),
        ("Obs. de Educacion y Formacion Profesional",      "241", "114.4", "0.475",
         "Alta demanda editorial — pilar de Ciencias de la Educacion", ACCENT),
        ("Obs. de Derecho y Justicia Social",              "127", " 64.8", "0.510",
         "10 articulos Q1 — eficiencia de impacto notable", ACCENT),
        ("Obs. de Desarrollo Social y Politicas Publicas", " 94", " 50.5", "0.537",
         "Nicho estrategico: consultoria y evaluacion de politica publica", LGRAY),
    ]

    y   = Inches(1.2)
    h_c = Inches(1.04)
    gap = Inches(0.08)
    col_w = W - Inches(0.8)

    for nombre, arts, imp, prom, desc, clr in obs:
        rect(sld, Inches(0.4), y, col_w, h_c, PANEL)
        rect(sld, Inches(0.4), y, Inches(0.08), h_c, clr)

        txt(sld, nombre, Inches(0.6), y + Inches(0.06),
            Inches(7.2), Inches(0.44), size=20, bold=True, color=WHITE)
        txt(sld, desc, Inches(0.6), y + Inches(0.54),
            Inches(7.2), Inches(0.42), size=14, italic=True, color=LGRAY)

        # Metricas
        for val, label, mx, mc in [
            (arts,  "Arts.",    Inches(8.3),  LGRAY),
            (imp,   "Impacto",  Inches(9.8),  GOLD),
            (prom,  "Prom.",    Inches(11.3), clr),
        ]:
            txt(sld, label, mx, y + Inches(0.06), Inches(1.4), Inches(0.34),
                size=11, color=LGRAY, align=PP_ALIGN.CENTER)
            txt(sld, val,   mx, y + Inches(0.42), Inches(1.4), Inches(0.52),
                size=20, bold=True, color=mc, align=PP_ALIGN.CENTER)

        y += h_c + gap

    footer(sld, 6)
    return sld


# ── S7: CALIDAD vs CANTIDAD ───────────────────────────────────────────────────
def s07_cuartiles(prs):
    sld = add_slide(prs)
    title_bar(sld, "Calidad vs. Cantidad: El Analisis de Cuartiles",
              "El volumen no equivale a impacto — los cuartiles revelan la verdad")

    img_bottom = insert_img(sld, F["cuartiles"], img_w=Inches(9.5), top=Inches(1.12))

    bullets = [
        ("El cluster con mayor impacto ponderado (0.639): Ingenieria Cuantitativa II "
         " con Q1=33 y Q2=31 sobre 366 articulos", GOLD),
        ("Obs. Multidisciplinario lidera en Q1 de la Via B: 41 articulos Q1 + 23 Q2", GREEN),
        ("Quimicas/Ambiental: alto volumen pero domina Q3/Q4 — necesita estrategia de upgrade", ACCENT),
        ("La metrica que guia la asignacion: impacto ponderado, no contador de papers", WHITE),
    ]
    by = img_bottom + Inches(0.18)
    bh = H - Inches(0.35) - by
    bullet_box(sld, bullets, Inches(0.5), by, W - Inches(1.0), bh,
               size=18, spacing=Pt(6), marker="->  ")

    footer(sld, 7)
    return sld


# ── S8: MASA CRITICA IPPC ─────────────────────────────────────────────────────
def s08_ippc(prs):
    sld = add_slide(prs)
    title_bar(sld, "Masa Critica IPPC: Con Quien Activamos los Centros",
              "Investigadores clasificados por percentil de productividad individual")

    img_bottom = insert_img(sld, F["ippc"], img_w=Inches(9.0), top=Inches(1.12))

    bullets = [
        ("Elite (>=P90): directores naturales de los nuevos centros — nucleo irremplazable", GOLD),
        ("Consolidados (P50-P89): columna vertebral operativa de cada Centro/Observatorio", ACCENT),
        ("FCQS y FCA validan su rol como matrices experimentales con investigadores Elite", GREEN),
        ("Proximo paso Fase C: cruzar IPPC Elite x cluster tematico -> asignacion individual", WHITE),
    ]
    by = img_bottom + Inches(0.18)
    bh = H - Inches(0.35) - by
    bullet_box(sld, bullets, Inches(0.5), by, W - Inches(1.0), bh,
               size=18, spacing=Pt(6), marker="->  ")

    footer(sld, 8)
    return sld


# ── S9: TOP LINEAS ────────────────────────────────────────────────────────────
def s09_lineas(prs):
    sld = add_slide(prs)
    title_bar(sld, "Top Lineas de I+D por Tipo de Investigacion",
              "Sobre estas lineas productivas nacen los reglamentos — no hay genericos")

    img_bottom = insert_img(sld, F["lineas"], img_w=Inches(9.5), top=Inches(1.12))

    bullets = [
        ("Via A experimental: Ambiente, Sistemas Agroalimentarios, Salud, Ingenieria", ACCENT),
        ("Via B social: Economia/Empresa, Procesos Educativos, Justicia y Gobernabilidad", GREEN),
        ("3-5 lineas formales por Centro/Observatorio — con produccion demostrada", GOLD),
        ("Base tecnica para el Reglamento General (Fase D, Abril 2026)", WHITE),
    ]
    by = img_bottom + Inches(0.18)
    bh = H - Inches(0.35) - by
    bullet_box(sld, bullets, Inches(0.5), by, W - Inches(1.0), bh,
               size=18, spacing=Pt(6), marker="->  ")

    footer(sld, 9)
    return sld


# ── S10: ROADMAP & DECISION ───────────────────────────────────────────────────
def s10_roadmap(prs):
    sld = add_slide(prs)
    title_bar(sld, "Roadmap & Decisión",
              "Arquitectura validada. Siguiente paso: aprobar e iniciar")

    # Fases completadas
    y = Inches(1.2)
    rect(sld, Inches(0.4), y, W - Inches(0.8), Inches(0.5), GREEN)
    txt(sld, "COMPLETADO -- Fase A: Diagnostico (Feb 2026)  |  "
             "Fase B: Analisis completo (Mar 2026)",
        Inches(0.6), y + Inches(0.07), W - Inches(1.2), Inches(0.38),
        size=18, bold=True, color=BG)
    y += Inches(0.58)

    # Fase C y D en dos columnas
    col_w = Inches(5.9)
    for cx, header, clr, items in [
        (Inches(0.4), "FASE C — Propuesta  (Mar-Abr 2026)", ACCENT, [
            "C.1  Cruzar IPPC Elite con clusters tematicos",
            "C.2  Propuestas tecnicas por Centro/Observatorio",
            "C.3  Definicion formal de lineas de investigacion",
            "C.4  Presentacion a Consejos de Facultad",
        ]),
        (Inches(7.0), "FASE D — Reglamentacion  (Abr 2026)", GOLD, [
            "D.1  Revision normativa: LOES, CES, Estatuto",
            "D.2  Redaccion del Reglamento General",
            "D.3  Aprobacion por Consejo Universitario",
            "",
        ]),
    ]:
        rect(sld, cx, y, col_w, Inches(0.45), clr)
        txt(sld, header, cx + Inches(0.12), y + Inches(0.06),
            col_w, Inches(0.36), size=17, bold=True, color=BG)
        yi = y + Inches(0.5)
        for item in items:
            if item:
                rect(sld, cx, yi, col_w, Inches(0.58), PANEL)
                txt(sld, item, cx + Inches(0.12), yi + Inches(0.09),
                    col_w - Inches(0.2), Inches(0.44), size=18, color=LGRAY)
            yi += Inches(0.63)

    # CTA
    cta_y = Inches(5.55)
    rect(sld, Inches(0.4), cta_y, W - Inches(0.8), Inches(1.5), PANEL)
    rect(sld, Inches(0.4), cta_y, Inches(0.1), Inches(1.5), GOLD)
    txt(sld, "DECISION REQUERIDA HOY",
        Inches(0.65), cta_y + Inches(0.1), W - Inches(1.2), Inches(0.45),
        size=22, bold=True, color=GOLD)
    txt(sld,
        "Arquitectura validada. Aprobar e iniciar asignacion de IPPC Elite a Centros "
        "y redaccion del Reglamento General para aprobacion por Consejo Universitario.",
        Inches(0.65), cta_y + Inches(0.6), W - Inches(1.2), Inches(0.78),
        size=20, color=WHITE)

    footer(sld, 10)
    return sld


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    print("=" * 55)
    print("  Generando Presentacion_Centros_UTMACH_Final.pptx")
    print("=" * 55)

    prs = new_prs()

    slides = [
        ("01 Titulo",                s01_titulo),
        ("02 El Problema",           s02_problema),
        ("03 ADN Cientifico",        s03_adn),
        ("04 Sistema Dos Vias",      s04_dos_vias),
        ("05 Via A - Centros",       s05_via_a),
        ("06 Via B - Observatorios", s06_via_b),
        ("07 Calidad vs Cantidad",   s07_cuartiles),
        ("08 Masa Critica IPPC",     s08_ippc),
        ("09 Top Lineas I+D",        s09_lineas),
        ("10 Roadmap y Decision",    s10_roadmap),
    ]

    for label, fn in slides:
        print(f"  [{label}]")
        fn(prs)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)

    print()
    print(f"[OK]  {OUTPUT}")
    print()
    print("  Figuras:")
    for key, path in F.items():
        ok = "[OK]" if os.path.exists(path) else "[FALTA]"
        print(f"    {ok}  {os.path.basename(path)}")

    print()
    print(f"  Total slides: {len(prs.slides)}")
    print("=" * 55)


if __name__ == "__main__":
    main()
