"""
==========================================================================
GENERADOR DE PRESENTACIÓN EJECUTIVA (.PPTX)
Centros de Investigación UTMACH - Sistema de Dos Vías
==========================================================================
"""
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

OUT_DIR = r"c:\Users\andre\Mi unidad\DIRECCIÓN DE INVESTIGACIÓN\Centros de Investigación\02_analisis\resultados"
FIG_DIR = os.path.join(OUT_DIR, "figuras")
PPTX_PATH = os.path.join(OUT_DIR, "Presentacion_DIDI_UTMACH_v2.pptx")

# Inicializar presentación
prs = Presentation()

# ==========================================
# RUTINAS AUXILIARES
# ==========================================
def add_title_slide(prs, title, subtitle):
    # Layout 0 es típicamente la diapositiva de título
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    
    title_box.text = title
    subtitle_box.text = subtitle
    return slide

def add_bullet_slide(prs, title, bullets):
    # Layout 1 es título y contenido
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_box.text = title
    tf = body_shape.text_frame
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
    # Ajustar fuente
    for p in tf.paragraphs:
        p.font.size = Pt(22)
            
    return slide

def add_image_slide(prs, title, img_name, bullets, img_width=Inches(5.5), img_left=Inches(4)):
    # Layout 5: Solo título, perfecto para customizar con imágenes
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    
    # Agregar texto a la izquierda
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
        p.font.size = Pt(18)
        p.space_after = Pt(14)
    
    # Agregar imagen a la derecha
    img_path = os.path.join(FIG_DIR, img_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, img_left, Inches(1.8), width=img_width)
    else:
        print(f"⚠️ Imagen no encontrada: {img_name}")
        
    return slide

# ==========================================
# CONSTRUCCIÓN DE DIAPOSITIVAS
# ==========================================

print("Generando presentación: Presentacion_DIDI_UTMACH.pptx...")

# Slide 1: Título
add_title_slide(prs, 
    "De la Intuición a la Evidencia:\nArquitectura Data-Driven UTMACH", 
    "Propuesta de Sistema de Dos Vías\nCentros Científico-Experimentales & Observatorios Sociales\n\nCorpus: 2,919 Artículos\nMetodología: NLP Semántico + Embeddings"
)

# Slide 2: El Problema
add_bullet_slide(prs, "El Desafío Estratégico", [
    "❌ Crear centros por 'palabras repetidas' agrupa vocabulario, no conocimiento",
    "❌ Asignar cuotas por facultad prioriza la política sobre los datos reales",
    "❌ Tratar igual ciencias experimentales y sociales diluye presupuestos",
    "⚠️ El Resultado Antiguo: Laboratorios pagando proyectos de escritorio y Observatorios sin software.",
    "La directriz del Director (Reg. 007): Cero prioridades a priori. Los datos mandan."
])

# Slide 3: Metodología
add_bullet_slide(prs, "Innovación Metodológica (NLP)", [
    "1. Datos Crudos: 2,919 artículos procesados.",
    "2. Clasificador Semántico: Identificación de vocabulario experimental vs social.",
    "3. K-Means Clustering: Agrupación algorítmica por contexto semántico (K=20 óptimo).",
    "4. Impacto Ponderado: Q1 = 1.0, Q2 = 0.9, Q3 = 0.8, Q4 = 0.7.",
    "✅ Conclusión: Se agrupa por la CIENCIA que hacen, no por la facultad a la que pertenecen."
])

# Slide 4: El Gran Descubrimiento (Clasificación EPS)
add_image_slide(prs, "¿Qué investiga realmente la UTMACH?", "01_clasificacion_produccion.png", [
    "El 47.0% es Ciencia Experimental (Trabajo de laboratorio, campo, ensayo clínico).",
    "El 41.4% es Ciencia Social / Humanística (Trabajo teórico, documental, social).",
    "Implicación:",
    "La UTMACH tiene DOS ADNs científicos distintos. Financiarlos igual es un error técnico.",
    "Solución: Arquitectura diferenciada en Dos Vías."
])

# Slide 5: El Sistema de Dos Vías
add_bullet_slide(prs, "El Nuevo Paradigma de Inversión", [
    "VÍA A: Centros Científico-Experimentales",
    " • Base: 1,371 artículos",
    " • Inversión: Laboratorios, reactivos, equipos de alto costo",
    " • Horizonte: Largo plazo",
    "",
    "VÍA B: Observatorios Sociales y Humanísticos",
    " • Base: 1,209 artículos",
    " • Inversión: Software, bases de datos digitales, encuestadores",
    " • Gestión ágil e impacto en políticas públicas."
])

# Slide 6: Vía A - Centros
add_image_slide(prs, "VÍA A: Los 4 Macro-Centros Experimentales", "02_impacto_por_centro.png", [
    "El impacto ponderado define la jerarquía:",
    "1. Ciencias Químicas y Ambientales (278 arts)",
    "2. Sist. Cuantitativos y Aplicados (284 arts)",
    "3. Salud Integral y Biociencias (265 arts)",
    "4. Agroalimentaria y Sostenibilidad (178 arts)",
    " ",
    "FCQS ancla dos de los centros de mayor peso.",
    "FCA ancla competitivamente el Agroalimentario."
])

# Slide 7: Vía B - Observatorios
add_bullet_slide(prs, "VÍA B: Los 4 Observatorios Sociales", [
    "📊 Observatorio de Economía, Empresa e Innovación",
    " • Concentra 547 artículos (El más grande de la institución).",
    "📚 Observatorio de Educación y Formación Profesional",
    " • 241 artículos, altísima demanda de publicación.",
    "⚖️ Observatorio de Derecho y Justicia Social",
    " • Mayor eficiencia: Alta concentración de papers Q1.",
    "🤝 Observatorio Multidisciplinario de Ciencias y Sociedad",
    " • Intersección de psicología, comunicación y ciencias políticas."
])

# Slide 8: Calidad vs Cantidad (Cuartiles)
add_image_slide(prs, "Calidad vs Cantidad (Análisis Q1-Q4)", "03_cuartiles_por_centro.png", [
    "Volumen NO es Impacto:",
    "• C. Cuantitativos / Experimental Gral tiene volumen, pero lidera en Q1/Q2.",
    "• C. Química/Ambiental tiene volumen pero domina Q3/Q4. (Requiere estrategia upgrade).",
    "• Obs. de Economía tiene 40+ artículos Q1.",
    " ",
    "La meta: Invertir en los campos que ya saben publicar en alto impacto."
], img_width=Inches(6.0), img_left=Inches(3.8))

# Slide 9: Masa Crítica IPPC
add_image_slide(prs, "Líderes de la Transformación (IPPC)", "05_masa_critica_ippc.png", [
    "Los Centros no pueden estar vacíos.",
    "• Investigadores Élite (Clúster A): Serán los directores naturales.",
    "• Investigadores Consolidados (Clúster B): La fuerza operativa.",
    "• FCQS es el motor indiscutible de investigadores Élite para la Vía A.",
    " ",
    "Siguiente fase: Cruzar estos rangos IPPC con la matriz de 8 Centros/Observatorios."
])

# Slide 10: Top Líneas Estratégicas
add_image_slide(prs, "Líneas de Investigación Fundacionales", "04_top_lineas_por_tipo.png", [
    "Cada centro nacerá con líneas pre-probadas, no con burocracia.",
    "Vía A liderada por:",
    " • Ambiente y conservación",
    " • Producción agroalimentaria",
    "Vía B liderada por:",
    " • Desarrollo económico",
    " • Educación e Innovación"
], img_width=Inches(6.0), img_left=Inches(3.8))

# Slide 11: Dashboard Integral
add_image_slide(prs, "Dashboard Ejecutivo: Arquitectura Científica", "06_dashboard_ejecutivo.png", [
    "Visión Panorámica de la Universidad:",
    " ",
    "• 2,919 Artículos analizados.",
    "• 80% del impacto institucional cubierto",
    "  en una arquitectura limpia y auditable.",
    " ",
    "• Revisiones Bibliográficas (339 arts) son la",
    "  base del estado del arte transversal."
], img_width=Inches(5.0), img_left=Inches(4.5))

# Slide 12: Roadmap - Siguientes Pasos
add_bullet_slide(prs, "Roadmap de Implementación: De Datos a Realidad", [
    "FASES A y B: COMPLETADAS ✅ (Este reporte)",
    "",
    "FASE C: Propuesta Estructural (Abril 2026)",
    " 1. Cruce matriz IPPC Élite vs Clusters.",
    " 2. Aprobación técnica de los 4 Centros y 4 Observatorios.",
    " 3. Presentación a Consejos de Facultad (FCQS / FCA).",
    "",
    "FASE D: Reglamentación",
    " 1. Aprobación por Consejo Universitario del reglamento del Sistema Dos Vías."
])

# Slide 13: Cierre y Decisión
add_bullet_slide(prs, "Conclusión y Decisión Estratégica", [
    "La UTMACH NO necesita diseñar centros de investigación. Necesita institucionalizar los que ya existen implícitamente en su producción.",
    "",
    "1. La evidencia dicta separar presupuestos (Experimental vs. Social).",
    "2. La masa crítica IPPC existe y concentra el impacto Q1.",
    "",
    "¿Solicitud al Director DIDI?:",
    "• Aprobación de la 'Arquitectura de Dos Vías'.",
    "• Luz verde para iniciar fase reglamentaria Fase C."
])

# Guardar
prs.save(PPTX_PATH)
print(f"✅ ¡Presentación PPTX generada exitosamente en:\n{PPTX_PATH}")
