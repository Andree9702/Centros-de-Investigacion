"""
generate_pptx.py
Genera la presentación PPTX: "Arquitectura Institucional de Investigación UTMACH"
Sistema de Dos Vías — Centros Científico-Experimentales & Observatorios Sociales
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── RUTAS ───────────────────────────────────────────────────────────────────
BASE_DIR   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGURAS    = os.path.join(BASE_DIR, "resultados", "figuras")
OUTPUT     = os.path.join(BASE_DIR, "resultados", "Presentacion_DIDI_UTMACH.pptx")

FIGS = {
    4:  os.path.join(FIGURAS, "01_clasificacion_produccion.png"),
    6:  os.path.join(FIGURAS, "02_impacto_por_centro.png"),
    8:  os.path.join(FIGURAS, "03_cuartiles_por_centro.png"),
    9:  os.path.join(FIGURAS, "05_masa_critica_ippc.png"),
    10: os.path.join(FIGURAS, "04_top_lineas_por_tipo.png"),
    11: os.path.join(FIGURAS, "06_dashboard_ejecutivo.png"),
}

# ─── PALETA DE COLORES ────────────────────────────────────────────────────────
C_AZUL_OSC  = RGBColor(0x0D, 0x1B, 0x2A)   # Fondo oscuro principal
C_AZUL_MED  = RGBColor(0x1A, 0x3A, 0x5C)   # Fondo secundario / cajas
C_AZUL_ACE  = RGBColor(0x00, 0x7A, 0xC3)   # Acento azul vibrante
C_VERDE_ACE = RGBColor(0x00, 0xB4, 0x8A)   # Acento verde (Vía B)
C_ROJO_ACE  = RGBColor(0xE0, 0x3B, 0x3B)   # Acento rojo (alertas)
C_DORADO    = RGBColor(0xF5, 0xA6, 0x23)   # Acento dorado (highlights)
C_BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRIS_CLAR = RGBColor(0xCF, 0xD8, 0xE0)
C_GRIS_TEXT = RGBColor(0xA0, 0xB0, 0xBE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]          # completamente en blanco
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=None):
    """Rellena el fondo del slide con un color sólido."""
    color = color or C_AZUL_OSC
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, word_wrap=True):
    color = color or C_BLANCO
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_size=16, bold=False,
                  color=None, align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    color = color or C_BLANCO
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_textbox_multi(slide, lines, x, y, w, h,
                      font_size=16, color=None, word_wrap=True):
    """Crea un textbox con múltiples párrafos desde una lista de (texto, kwargs)."""
    color = color or C_BLANCO
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            item = {"text": item}
        text       = item.get("text", "")
        fsize      = item.get("font_size", font_size)
        bold       = item.get("bold", False)
        italic     = item.get("italic", False)
        clr        = item.get("color", color)
        align      = item.get("align", PP_ALIGN.LEFT)
        sp_before  = item.get("space_before", Pt(5))

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = sp_before
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fsize)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = clr
    return txBox


def add_image_centered(slide, img_path, top, height, slide_w=SLIDE_W):
    """Inserta una imagen centrada horizontalmente con altura fija."""
    from pptx.util import Inches
    from PIL import Image as PILImage
    try:
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        w = height * aspect
        if w > slide_w * 0.95:
            w = slide_w * 0.95
            height = w / aspect
        left = (slide_w - w) / 2
        slide.shapes.add_picture(img_path, left, top, width=w, height=height)
    except Exception:
        # Si PIL no está disponible, insertar con tamaño fijo
        w = min(height * 1.6, slide_w * 0.9)
        left = (slide_w - w) / 2
        slide.shapes.add_picture(img_path, left, top, width=w, height=height)


def header_bar(slide, title, subtitle=None):
    """Barra de encabezado azul oscuro con título."""
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, C_AZUL_MED)
    # Línea acento
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_AZUL_ACE)
    add_text(slide, title,
             Inches(0.4), Inches(0.12),
             SLIDE_W - Inches(0.8), Inches(0.7),
             font_size=26, bold=True, color=C_BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.75),
                 SLIDE_W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=False, color=C_AZUL_ACE, align=PP_ALIGN.LEFT)


def footer_bar(slide, text="UTMACH · DIDI · Marzo 2026 · Análisis de 2.919 artículos"):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), C_AZUL_MED)
    add_text(slide, text,
             Inches(0.3), SLIDE_H - Inches(0.30),
             SLIDE_W - Inches(0.6), Inches(0.28),
             font_size=9, color=C_GRIS_TEXT, align=PP_ALIGN.LEFT)
    # Número de página no nativo — omitimos por simplicidad


def accent_box(slide, text, x, y, w, h, bg=None, txt_color=None, font_size=15, bold=True):
    bg = bg or C_AZUL_ACE
    txt_color = txt_color or C_BLANCO
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, text, x + Inches(0.12), y + Inches(0.05),
             w - Inches(0.24), h - Inches(0.1),
             font_size=font_size, bold=bold, color=txt_color,
             align=PP_ALIGN.LEFT, word_wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_titulo(prs):
    """SLIDE 1 — Título y enganche"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)

    # Franja izquierda decorativa
    add_rect(sld, 0, 0, Inches(0.35), SLIDE_H, C_AZUL_ACE)

    # Título principal
    add_text(sld, "De la Intuición a la Evidencia:",
             Inches(0.6), Inches(0.9),
             Inches(12.0), Inches(1.1),
             font_size=36, bold=True, color=C_BLANCO)
    add_text(sld, "Una Arquitectura de Investigación Data-Driven para la UTMACH",
             Inches(0.6), Inches(1.9),
             Inches(12.0), Inches(0.9),
             font_size=26, bold=False, color=C_AZUL_ACE)

    # Separador
    add_rect(sld, Inches(0.6), Inches(2.85), Inches(10.0), Inches(0.04), C_DORADO)

    # Subtítulo
    add_text(sld, "Propuesta de Sistema de Dos Vías —  "
                  "Centros Científico-Experimentales  &  Observatorios Sociales/Humanísticos",
             Inches(0.6), Inches(3.0),
             Inches(12.0), Inches(0.8),
             font_size=18, bold=False, color=C_GRIS_CLAR)

    # Credenciales
    creds = [
        ("Corpus analizado:", C_DORADO, True),
        ("2.919 artículos científicos — producción histórica UTMACH", C_BLANCO, False),
        ("Metodología:", C_DORADO, True),
        ("NLP semántico + Embeddings + K-Means clustering", C_BLANCO, False),
        ("Equipo:", C_DORADO, True),
        ("PhD. Ivan Ramirez (DIDI)  ·  MSc. Andreé Vitonera  ·  MSc. Luiggi Solano", C_BLANCO, False),
        ("Fecha:", C_DORADO, True),
        ("Marzo 2026", C_BLANCO, False),
    ]
    y0 = Inches(3.95)
    dy = Inches(0.38)
    for i, (txt, clr, bold) in enumerate(creds):
        add_text(sld, txt,
                 Inches(0.6) if bold else Inches(2.5),
                 y0 + (i // 2) * dy,
                 Inches(2.0) if bold else Inches(9.0),
                 Inches(0.36),
                 font_size=14, bold=bold, color=clr)

    # Cita de enganche
    add_rect(sld, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.75), C_AZUL_MED)
    add_text(sld, '"Lo que van a ver hoy no es una propuesta de centros. '
                  'Es el resultado de dejar que los datos de la propia UTMACH hablen por sí mismos."',
             Inches(0.75), Inches(6.57), Inches(11.7), Inches(0.6),
             font_size=13, italic=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_02_problema(prs):
    """SLIDE 2 — El Problema"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Por qué el Método Tradicional Falla",
               "El costo estratégico de crear centros sin datos")

    # Columna izquierda — El modelo antiguo
    add_rect(sld, Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.42), C_ROJO_ACE)
    add_text(sld, "❌  EL MODELO ANTIGUO",
             Inches(0.5), Inches(1.28), Inches(5.6), Inches(0.38),
             font_size=14, bold=True, color=C_BLANCO)

    errores = [
        "Crear centros por «palabras repetidas en títulos»\n→ agrupa vocabulario, no conocimiento real",
        "Asignar cuotas por facultad\n→ sesgos políticos sobre evidencia científica",
        "Tratar igual la investigación experimental\ny la documental → presupuestos incorrectos",
        "Priorizar facultades a priori\n→ sesga el análisis antes de comenzar (Registro 007)",
    ]
    y = Inches(1.75)
    for e in errores:
        add_text(sld, e, Inches(0.55), y, Inches(5.55), Inches(0.65),
                 font_size=13, color=C_GRIS_CLAR)
        y += Inches(0.75)

    # Columna derecha — El costo
    add_rect(sld, Inches(6.8), Inches(1.25), Inches(6.1), Inches(0.42), C_DORADO)
    add_text(sld, "⚠  EL COSTO ESTRATÉGICO",
             Inches(6.9), Inches(1.28), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=C_AZUL_OSC)

    costos = [
        "Laboratorios financiando proyectos de escritorio",
        "Observatorios sociales pagando reactivos que nunca usarán",
        "Investigadores élite en centros fuera de su línea de impacto",
        "Recursos diluidos · impacto institucional mínimo",
    ]
    y = Inches(1.75)
    for c in costos:
        add_rect(sld, Inches(6.85), y, Inches(5.9), Inches(0.6), C_AZUL_MED)
        add_text(sld, c, Inches(7.0), y + Inches(0.08), Inches(5.6), Inches(0.52),
                 font_size=13, color=C_GRIS_CLAR)
        y += Inches(0.75)

    # Separador vertical
    add_rect(sld, Inches(6.5), Inches(1.25), Inches(0.04), Inches(3.5), C_AZUL_ACE)

    # Frase de cierre
    add_rect(sld, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.9), C_AZUL_MED)
    add_text(sld, '"Un centro diseñado por cuota política no produce ciencia. Produce burocracia."',
             Inches(0.6), Inches(6.38), Inches(12.1), Inches(0.65),
             font_size=16, italic=True, bold=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_03_metodologia(prs):
    """SLIDE 3 — Innovación Metodológica"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Innovación Metodológica",
               "NLP + Semántica + Ponderación por Impacto — cómo lo resolvimos")

    # Pipeline de 4 pasos
    pasos = [
        ("1", "DATOS BRUTOS", "2.919 artículos\nTítulos · Abstracts\nLíneas · Campos"),
        ("2", "CLASIFICADOR NLP", "Léxico semántico\nExp. vs Social\nvs Bibliográfica"),
        ("3", "CLUSTERING K-MEANS", "Silhouette Score\nK=20 óptimo\nTF-IDF + bigramas"),
        ("4", "PONDERACIÓN Q1-Q4", "Q1=1.0 / Q2=0.9\nQ3=0.8 / Q4=0.7\nLatindex=0.2"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.9)
    gap   = Inches(0.35)
    total_w = 4 * box_w + 3 * gap
    x0 = (SLIDE_W - total_w) / 2
    y0 = Inches(1.4)

    for i, (num, titulo, cuerpo) in enumerate(pasos):
        x = x0 + i * (box_w + gap)
        clr = [C_AZUL_ACE, C_VERDE_ACE, C_DORADO, C_ROJO_ACE][i]
        add_rect(sld, x, y0, box_w, box_h, clr)
        add_text(sld, num, x + Inches(0.1), y0 + Inches(0.05), Inches(0.5), Inches(0.45),
                 font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        add_text(sld, titulo, x + Inches(0.1), y0 + Inches(0.45), box_w - Inches(0.2), Inches(0.45),
                 font_size=13, bold=True, color=C_BLANCO)
        add_text(sld, cuerpo, x + Inches(0.1), y0 + Inches(0.9), box_w - Inches(0.2), Inches(0.95),
                 font_size=11, color=C_BLANCO)

        # Flecha entre cajas
        if i < 3:
            ax = x + box_w + Inches(0.05)
            add_text(sld, "→", ax, y0 + Inches(0.65), gap - Inches(0.05), Inches(0.5),
                     font_size=22, bold=True, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    # Decisiones clave
    add_rect(sld, Inches(0.4), Inches(3.55), Inches(12.5), Inches(0.04), C_AZUL_ACE)

    bullets = [
        ("✅  Eliminación de prioridades a priori — los datos determinan las fortalezas  (Registro 007)", C_VERDE_ACE),
        ("✅  Clasificación Experimental vs. Social antes del clustering → inversión diferenciada  (Registro 010)", C_VERDE_ACE),
        ("✅  Ponderación aprobada: «No es lo mismo 100 artículos Latindex que 50 Q1»  (Registro 008)", C_DORADO),
        ("✅  87.3% de artículos sin cuartil Scopus/WoS → la ponderación es aún más crítica  (Registro 009)", C_DORADO),
    ]
    y = Inches(3.75)
    for txt, clr in bullets:
        add_text(sld, txt, Inches(0.6), y, Inches(12.1), Inches(0.45),
                 font_size=13, color=clr)
        y += Inches(0.5)

    footer_bar(sld)
    return sld


def slide_04_epistemologia(prs):
    """SLIDE 4 — El Gran Descubrimiento + Figura 01"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "El Gran Descubrimiento: La Epistemología de la UTMACH",
               "¿Qué tipo de universidad es la UTMACH? Los datos responden")

    # Figura a la izquierda
    img = FIGS[4]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(4.0),
                           slide_w=Inches(6.5))
        fig_x_end = Inches(6.8)
    else:
        add_rect(sld, Inches(0.4), Inches(1.2), Inches(6.0), Inches(4.0), C_AZUL_MED)
        add_text(sld, "[01_clasificacion_produccion.png]",
                 Inches(0.5), Inches(2.8), Inches(5.8), Inches(0.6),
                 font_size=12, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)
        fig_x_end = Inches(6.8)

    # Datos y análisis a la derecha
    datos = [
        ("47.0%", "EXPERIMENTAL", "1.371 artículos\nTrabajo de laboratorio, campo o clínico", C_AZUL_ACE),
        ("41.4%", "CS / HUMANÍSTICAS", "1.209 artículos\nTeóricos, documentales, socioeconómicos", C_VERDE_ACE),
        ("11.6%", "REVISIONES BIBL.", "339 artículos\nMeta-análisis, estados del arte", C_DORADO),
    ]
    y = Inches(1.3)
    for pct, label, desc, clr in datos:
        add_rect(sld, fig_x_end, y, Inches(1.0), Inches(1.1), clr)
        add_text(sld, pct, fig_x_end + Inches(0.05), y + Inches(0.15),
                 Inches(0.9), Inches(0.6), font_size=22, bold=True, color=C_BLANCO,
                 align=PP_ALIGN.CENTER)
        add_rect(sld, fig_x_end + Inches(1.05), y, Inches(5.0), Inches(1.1), C_AZUL_MED)
        add_text(sld, label, fig_x_end + Inches(1.15), y + Inches(0.05),
                 Inches(4.8), Inches(0.35), font_size=13, bold=True, color=clr)
        add_text(sld, desc, fig_x_end + Inches(1.15), y + Inches(0.4),
                 Inches(4.8), Inches(0.6), font_size=11, color=C_GRIS_CLAR)
        y += Inches(1.25)

    # Implicación estratégica
    add_rect(sld, fig_x_end, Inches(5.1), Inches(6.1), Inches(1.0), C_ROJO_ACE)
    add_text(sld,
             "La UTMACH tiene DOS ADNs científicos de igual peso.\n"
             "Financiarlos con la misma lógica presupuestaria es un error técnico y financiero.",
             fig_x_end + Inches(0.12), Inches(5.18), Inches(5.85), Inches(0.9),
             font_size=12, bold=True, color=C_BLANCO)

    footer_bar(sld)
    return sld


def slide_05_paradigma(prs):
    """SLIDE 5 — El Nuevo Paradigma"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "El Nuevo Paradigma: El Sistema de Dos Vías",
               "Una arquitectura financiera inteligente para dos ciencias diferentes")

    col_w = Inches(5.9)
    col_h = Inches(5.1)
    y0 = Inches(1.3)

    # VÍA A
    add_rect(sld, Inches(0.4), y0, col_w, Inches(0.5), C_AZUL_ACE)
    add_text(sld, "🔬  VÍA A — CENTROS CIENTÍFICO-EXPERIMENTALES",
             Inches(0.5), y0 + Inches(0.06), col_w - Inches(0.2), Inches(0.42),
             font_size=13, bold=True, color=C_BLANCO)
    add_rect(sld, Inches(0.4), y0 + Inches(0.5), col_w, col_h - Inches(0.5), C_AZUL_MED)

    via_a = [
        ("Base científica", "1.371 artículos (47.0%)"),
        ("Infraestructura", "Laboratorios · reactivos · equipos especializados"),
        ("Operación", "Protocolo rígido · bioseguridad · cadena de frío"),
        ("Presupuesto", "Alto / Inversión de capital a largo plazo"),
        ("Estructuras", "4 Macro-Centros de Investigación"),
    ]
    y = y0 + Inches(0.65)
    for campo, valor in via_a:
        add_text(sld, campo + ":", Inches(0.55), y, Inches(1.8), Inches(0.42),
                 font_size=12, bold=True, color=C_AZUL_ACE)
        add_text(sld, valor, Inches(2.35), y, Inches(3.85), Inches(0.42),
                 font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.7)

    # VÍA B
    x2 = Inches(7.0)
    add_rect(sld, x2, y0, col_w, Inches(0.5), C_VERDE_ACE)
    add_text(sld, "📊  VÍA B — OBSERVATORIOS SOCIALES/HUMANÍSTICOS",
             x2 + Inches(0.1), y0 + Inches(0.06), col_w - Inches(0.2), Inches(0.42),
             font_size=13, bold=True, color=C_BLANCO)
    add_rect(sld, x2, y0 + Inches(0.5), col_w, col_h - Inches(0.5), C_AZUL_MED)

    via_b = [
        ("Base científica", "1.209 artículos (41.4%)"),
        ("Infraestructura", "Software · bases de datos · plataformas digitales"),
        ("Operación", "Gestión ágil · trabajo remoto · encuestadores de campo"),
        ("Presupuesto", "Moderado / Operación flexible y ágil"),
        ("Estructuras", "4 Macro-Observatorios Académicos"),
    ]
    y = y0 + Inches(0.65)
    for campo, valor in via_b:
        add_text(sld, campo + ":", x2 + Inches(0.15), y, Inches(1.8), Inches(0.42),
                 font_size=12, bold=True, color=C_VERDE_ACE)
        add_text(sld, valor, x2 + Inches(1.95), y, Inches(3.85), Inches(0.42),
                 font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.7)

    # Separador VS
    add_rect(sld, Inches(6.5), y0, Inches(0.06), col_h, C_DORADO)
    add_text(sld, "VS", Inches(6.2), y0 + Inches(2.1), Inches(0.7), Inches(0.6),
             font_size=20, bold=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    # Frase de cierre
    add_rect(sld, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.65), C_AZUL_MED)
    add_text(sld, "Dos vías, una sola estrategia: asignar recursos donde existe capacidad instalada real.",
             Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.5),
             font_size=14, italic=True, bold=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_06_via_a(prs):
    """SLIDE 6 — VÍA A: Los 4 Centros + Figura 02"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "VÍA A — Los 4 Centros Científico-Experimentales",
               "Masa crítica e impacto ponderado demostrado")

    # Figura arriba
    img = FIGS[6]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(3.0))
    else:
        add_rect(sld, Inches(1.5), Inches(1.2), Inches(10.0), Inches(3.0), C_AZUL_MED)
        add_text(sld, "[02_impacto_por_centro.png]",
                 Inches(2.0), Inches(2.5), Inches(9.0), Inches(0.6),
                 font_size=12, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    # Tabla de centros
    centros = [
        ("Ciencias Químicas y Ambientales", "278",  "138.2", "0.497", C_AZUL_ACE),
        ("Experimental General",            "284",  "137.2", "0.483", C_VERDE_ACE),
        ("Salud Integral y Biociencias",    "265",  "126.6", "0.478", C_DORADO),
        ("Agroalimentaria y Sostenibilidad","178",  " 91.2", "0.512", C_ROJO_ACE),
    ]
    headers = ["CENTRO", "ARTÍCULOS", "IMPACTO PON.", "PROM."]
    col_xs = [Inches(0.4), Inches(6.4), Inches(8.6), Inches(10.5)]
    col_ws = [Inches(5.9), Inches(2.1), Inches(1.8), Inches(1.8)]
    y_hdr = Inches(4.4)

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        add_rect(sld, cx, y_hdr, cw - Inches(0.05), Inches(0.35), C_AZUL_ACE)
        add_text(sld, hdr, cx + Inches(0.05), y_hdr + Inches(0.04),
                 cw - Inches(0.1), Inches(0.28), font_size=10, bold=True,
                 color=C_BLANCO, align=PP_ALIGN.CENTER)

    for i, (nombre, arts, impacto, prom, clr) in enumerate(centros):
        y_row = y_hdr + Inches(0.35) + i * Inches(0.48)
        bg = C_AZUL_MED if i % 2 == 0 else C_AZUL_OSC
        add_rect(sld, col_xs[0], y_row, Inches(12.5), Inches(0.46), bg)
        add_rect(sld, col_xs[0], y_row, Inches(0.07), Inches(0.46), clr)
        vals = [nombre, arts, impacto, prom]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
        for j, (val, cx, cw, al) in enumerate(zip(vals, col_xs, col_ws, aligns)):
            fc = C_DORADO if j == 2 else C_BLANCO
            fs = 11 if j == 0 else 13
            bold = j == 2
            add_text(sld, val, cx + Inches(0.1), y_row + Inches(0.07),
                     cw - Inches(0.15), Inches(0.38), font_size=fs,
                     bold=bold, color=fc, align=al)

    footer_bar(sld)
    return sld


def slide_07_via_b(prs):
    """SLIDE 7 — VÍA B: Los 4 Observatorios"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "VÍA B — Los 4 Observatorios Sociales/Humanísticos",
               "El mayor volumen de producción de la UTMACH — gestión ágil y digital")

    obs = [
        ("Economía, Empresa e Innovación",          "547", "269.6", "0.493",
         "El más grande de la institución (18.7% de toda la producción).", C_AZUL_ACE),
        ("Educación y Formación Profesional",       "241", "114.4", "0.475",
         "Alta demanda de publicación. Pilar de la FEDU.", C_VERDE_ACE),
        ("Derecho y Justicia Social",               "127", " 64.8", "0.510",
         "10 artículos Q1 — mayor eficiencia Q1 de la Vía B.", C_DORADO),
        ("Desarrollo Social y Políticas Públicas",  " 94", " 50.5", "0.537",
         "Nicho estratégico para consultoría y políticas públicas.", C_ROJO_ACE),
    ]

    y0 = Inches(1.4)
    card_h = Inches(1.15)
    gap    = Inches(0.12)
    for i, (nombre, arts, impacto, prom, desc, clr) in enumerate(obs):
        y = y0 + i * (card_h + gap)
        add_rect(sld, Inches(0.4), y, Inches(0.1), card_h, clr)
        add_rect(sld, Inches(0.5), y, Inches(12.4), card_h, C_AZUL_MED)

        add_text(sld, nombre, Inches(0.65), y + Inches(0.1),
                 Inches(7.0), Inches(0.45), font_size=15, bold=True, color=C_BLANCO)
        add_text(sld, desc, Inches(0.65), y + Inches(0.58),
                 Inches(7.0), Inches(0.45), font_size=11, color=C_GRIS_CLAR, italic=True)

        # Métricas
        mets = [("Artículos", arts, C_GRIS_CLAR), ("Impacto", impacto, C_DORADO), ("Prom.", prom, C_AZUL_ACE)]
        mx = Inches(8.0)
        for label, val, mc in mets:
            add_text(sld, label, mx, y + Inches(0.1), Inches(1.4), Inches(0.35),
                     font_size=10, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)
            add_text(sld, val, mx, y + Inches(0.42), Inches(1.4), Inches(0.55),
                     font_size=18, bold=True, color=mc, align=PP_ALIGN.CENTER)
            mx += Inches(1.55)

    # Nota operativa
    add_rect(sld, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.85), C_AZUL_MED)
    add_text(sld,
             "💡  Costo operativo reducido: sin laboratorios. Inversión en software analítico, "
             "licencias de bases de datos y encuestas de campo.\n"
             "🔗  Potencial de consultoría y transferencia de conocimiento al sector público y privado.",
             Inches(0.6), Inches(6.38), Inches(12.1), Inches(0.75),
             font_size=12, color=C_GRIS_CLAR)

    footer_bar(sld)
    return sld


def slide_08_cuartiles(prs):
    """SLIDE 8 — Calidad vs Cantidad + Figura 03"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Calidad vs. Cantidad: El Análisis de Cuartiles",
               "Más artículos no significa más impacto — los cuartiles revelan la verdad")

    img = FIGS[8]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(3.5))
    else:
        add_rect(sld, Inches(1.5), Inches(1.2), Inches(10.0), Inches(3.5), C_AZUL_MED)
        add_text(sld, "[03_cuartiles_por_centro.png]",
                 Inches(2.0), Inches(2.7), Inches(9.0), Inches(0.6),
                 font_size=12, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    insights = [
        ("🏆", "Clúster Experimental 6 lidera en calidad", C_DORADO,
         "Solo 179 artículos pero Q1=29 · Q2=30 · Prom. ponderado: 0.639 — el más alto de toda la Vía A"),
        ("📈", "Obs. Social Clúster 4 lidera en Vía B", C_VERDE_ACE,
         "41 artículos Q1 + 23 Q2 — mayor potencial de visibilidad internacional en la Vía B"),
        ("⚠", "C. Químicas/Ambiental requiere estrategia de upgrade", C_ROJO_ACE,
         "Alto volumen (278 art.) pero domina Q3/Q4 → focalizar en revistas Q1-Q2"),
    ]
    y = Inches(4.95)
    for icon, title, clr, desc in insights:
        add_rect(sld, Inches(0.4), y, Inches(12.5), Inches(0.6), C_AZUL_MED)
        add_text(sld, icon + "  " + title, Inches(0.55), y + Inches(0.08),
                 Inches(5.5), Inches(0.44), font_size=13, bold=True, color=clr)
        add_text(sld, desc, Inches(6.1), y + Inches(0.1),
                 Inches(6.7), Inches(0.44), font_size=11, color=C_GRIS_CLAR, italic=True)
        y += Inches(0.68)

    add_rect(sld, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.2), C_AZUL_MED)
    add_text(sld,
             "La métrica de éxito no es cuánto se publica. Es DÓNDE se publica. "
             "1 artículo Q1 = 5 artículos Q4 en impacto institucional.",
             Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.38),
             font_size=11, italic=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_09_ippc(prs):
    """SLIDE 9 — Masa Crítica IPPC + Figura 05"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Los Líderes de la Transformación: Masa Crítica IPPC",
               "¿Con quién llenamos estos centros? El talento identificado por datos")

    img = FIGS[9]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(3.3),
                           slide_w=Inches(7.5))

    # Leyenda clusters IPPC
    clusters = [
        ("ÉLITE  ≥ P90",         "Núcleo irremplazable. Directores naturales de los nuevos centros.", C_DORADO),
        ("CONSOLIDADOS  P50-P89","Columna vertebral. Investigadores con trayectoria sostenida.",     C_AZUL_ACE),
        ("EN DESARROLLO  P25-P49","Pipeline futuro. Candidatos a programas de mentoría.",             C_VERDE_ACE),
        ("SIN ACTIVIDAD  < P25", "Requieren plan de fortalecimiento o reorientación.",                C_GRIS_TEXT),
    ]
    y = Inches(1.3)
    x_leg = Inches(7.7)
    for label, desc, clr in clusters:
        add_rect(sld, x_leg, y, Inches(5.3), Inches(0.55), C_AZUL_MED)
        add_rect(sld, x_leg, y, Inches(0.12), Inches(0.55), clr)
        add_text(sld, label, x_leg + Inches(0.2), y + Inches(0.02),
                 Inches(5.0), Inches(0.28), font_size=12, bold=True, color=clr)
        add_text(sld, desc, x_leg + Inches(0.2), y + Inches(0.28),
                 Inches(5.0), Inches(0.25), font_size=10, color=C_GRIS_CLAR, italic=True)
        y += Inches(0.65)

    # Validaciones
    add_rect(sld, Inches(0.4), Inches(4.65), Inches(12.5), Inches(0.04), C_AZUL_ACE)

    validaciones = [
        ("🔬", "FCQS es pilar de VÍA A:", "Concentra investigadores élite en salud, química y biociencias.", C_AZUL_ACE),
        ("🌾", "FCA ancla la Vía A Agroalimentaria:", "Justificación técnica de la decisión piloto (Registro 003).", C_VERDE_ACE),
        ("📋", "Fase C — próximo paso:", "Cruzar IPPC élite con clusters temáticos → asignación individual por investigador.", C_DORADO),
    ]
    y = Inches(4.8)
    for icon, bold_txt, normal_txt, clr in validaciones:
        add_text(sld, icon + "  " + bold_txt, Inches(0.5), y,
                 Inches(4.5), Inches(0.42), font_size=13, bold=True, color=clr)
        add_text(sld, normal_txt, Inches(5.0), y,
                 Inches(7.9), Inches(0.42), font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.55)

    footer_bar(sld)
    return sld


def slide_10_lineas(prs):
    """SLIDE 10 — Líneas estratégicas + Figura 04"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Líneas Estratégicas por Tipo de Investigación",
               "Donde concentrar la inversión — respaldado por producción real")

    img = FIGS[10]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(3.6))
    else:
        add_rect(sld, Inches(1.5), Inches(1.2), Inches(10.0), Inches(3.6), C_AZUL_MED)
        add_text(sld, "[04_top_lineas_por_tipo.png]",
                 Inches(2.0), Inches(2.7), Inches(9.0), Inches(0.6),
                 font_size=12, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    add_rect(sld, Inches(0.4), Inches(5.0), Inches(5.9), Inches(0.35), C_AZUL_ACE)
    add_text(sld, "🔬  LÍNEAS VÍA A (Experimental)",
             Inches(0.5), Inches(5.04), Inches(5.7), Inches(0.28),
             font_size=12, bold=True, color=C_BLANCO)

    lineas_a = [
        "Ambiente y conservación",
        "Sistemas agroalimentarios sostenibles",
        "Prevención, promoción y cuidados de salud",
        "Manejo integral de entidades nosológicas",
        "Procesos químicos y materiales sostenibles",
    ]
    y = Inches(5.42)
    for l in lineas_a:
        add_text(sld, "▸  " + l, Inches(0.6), y, Inches(5.6), Inches(0.38),
                 font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.36)

    add_rect(sld, Inches(7.0), Inches(5.0), Inches(5.9), Inches(0.35), C_VERDE_ACE)
    add_text(sld, "📊  LÍNEAS VÍA B (Social/Humanística)",
             Inches(7.1), Inches(5.04), Inches(5.7), Inches(0.28),
             font_size=12, bold=True, color=C_BLANCO)

    lineas_b = [
        "Desarrollo económico y empresarial",
        "Procesos educativos y formación humana",
        "Justicia y gobernabilidad",
        "Justicia social, desarrollo humano y PP",
        "Comunicación digital, educación e innovación",
    ]
    y = Inches(5.42)
    for l in lineas_b:
        add_text(sld, "▸  " + l, Inches(7.15), y, Inches(5.6), Inches(0.38),
                 font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.36)

    add_rect(sld, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.15), C_AZUL_MED)
    add_text(sld, "Estas líneas son la base para los Reglamentos de los Centros (Fase D) — no genéricas, sino con producción demostrada.",
             Inches(0.6), Inches(7.12), Inches(12.1), Inches(0.3),
             font_size=10, italic=True, color=C_DORADO, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_11_dashboard(prs):
    """SLIDE 11 — Dashboard Ejecutivo + Figura 06"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Dashboard Ejecutivo: Panorama Completo",
               "El ecosistema de investigación UTMACH en una sola vista")

    img = FIGS[11]
    if os.path.exists(img):
        add_image_centered(sld, img, top=Inches(1.2), height=Inches(4.5))
    else:
        add_rect(sld, Inches(1.0), Inches(1.2), Inches(11.0), Inches(4.5), C_AZUL_MED)
        add_text(sld, "[06_dashboard_ejecutivo.png]",
                 Inches(2.0), Inches(3.0), Inches(9.0), Inches(0.6),
                 font_size=12, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    kpis = [
        ("2.919", "artículos analizados", C_AZUL_ACE),
        ("8",     "unidades de investigación (4+4)", C_VERDE_ACE),
        ("88.4%", "producción con asignación estructurada", C_DORADO),
        ("~1.120", "puntos de impacto ponderado total", C_ROJO_ACE),
    ]
    kpi_w = SLIDE_W / len(kpis) - Inches(0.2)
    x = Inches(0.4)
    for val, label, clr in kpis:
        add_rect(sld, x, Inches(5.95), kpi_w, Inches(1.15), C_AZUL_MED)
        add_rect(sld, x, Inches(5.95), kpi_w, Inches(0.08), clr)
        add_text(sld, val, x + Inches(0.1), Inches(6.1), kpi_w - Inches(0.2), Inches(0.55),
                 font_size=26, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_text(sld, label, x + Inches(0.1), Inches(6.65), kpi_w - Inches(0.2), Inches(0.38),
                 font_size=10, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)
        x += kpi_w + Inches(0.2)

    footer_bar(sld)
    return sld


def slide_12_roadmap(prs):
    """SLIDE 12 — Plan de Acción / Roadmap"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Plan de Acción — Roadmap: Fases C y D",
               "De los datos a la reglamentación — el camino ya trazado")

    # Fases completadas
    add_rect(sld, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.42), C_VERDE_ACE)
    add_text(sld, "✅  COMPLETADO — Fase A: Diagnóstico y Datos (Febrero 2026)  "
                  "·  Fase B: Análisis de Viabilidad (Marzo 2026)",
             Inches(0.55), Inches(1.35), Inches(12.1), Inches(0.35),
             font_size=12, bold=True, color=C_BLANCO)

    # Fase C
    fc_y = Inches(1.95)
    add_rect(sld, Inches(0.4), fc_y, Inches(5.9), Inches(0.42), C_AZUL_ACE)
    add_text(sld, "⏳  FASE C — Propuesta de Centros  (Marzo – Abril 2026)",
             Inches(0.5), fc_y + Inches(0.05), Inches(5.7), Inches(0.35),
             font_size=12, bold=True, color=C_BLANCO)
    fase_c = [
        "C.1  →  Cruzar investigadores élite IPPC con clusters temáticos",
        "C.2  →  Elaborar propuestas técnicas por cada Centro / Observatorio",
        "C.3  →  Definición formal de líneas de investigación por unidad",
        "C.4  →  Presentación a Consejos de Facultad",
    ]
    y = fc_y + Inches(0.5)
    for item in fase_c:
        add_rect(sld, Inches(0.4), y, Inches(5.9), Inches(0.54), C_AZUL_MED)
        add_text(sld, item, Inches(0.55), y + Inches(0.08),
                 Inches(5.65), Inches(0.4), font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.6)

    # Fase D
    fd_x = Inches(7.0)
    fd_y = Inches(1.95)
    add_rect(sld, fd_x, fd_y, Inches(5.9), Inches(0.42), C_DORADO)
    add_text(sld, "⏳  FASE D — Reglamentación  (Abril 2026)",
             fd_x + Inches(0.1), fd_y + Inches(0.05), Inches(5.7), Inches(0.35),
             font_size=12, bold=True, color=C_AZUL_OSC)
    fase_d = [
        "D.1  →  Revisión normativa: LOES, CES, Estatuto Universitario",
        "D.2  →  Redacción del Reglamento General de Centros y Observatorios",
        "D.3  →  Aprobación por Consejo Universitario",
    ]
    y = fd_y + Inches(0.5)
    for item in fase_d:
        add_rect(sld, fd_x, y, Inches(5.9), Inches(0.54), C_AZUL_MED)
        add_text(sld, item, fd_x + Inches(0.15), y + Inches(0.08),
                 Inches(5.65), Inches(0.4), font_size=12, color=C_GRIS_CLAR)
        y += Inches(0.6)

    # Separador vertical
    add_rect(sld, Inches(6.55), Inches(1.95), Inches(0.06), Inches(2.7), C_AZUL_ACE)

    # Decisión clave
    add_rect(sld, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.55), C_AZUL_MED)
    add_rect(sld, Inches(0.4), Inches(5.6), Inches(0.12), Inches(1.55), C_DORADO)
    add_text(sld, "DECISIÓN QUE SE REQUIERE HOY",
             Inches(0.7), Inches(5.68), Inches(12.0), Inches(0.38),
             font_size=14, bold=True, color=C_DORADO)
    add_text(sld,
             "Aprobación de la arquitectura de dos vías (4 Centros + 4 Observatorios).\n"
             "Esto activa de inmediato la Fase C: cruce IPPC élite × clusters temáticos\n"
             "y la presentación técnica a cada Facultad durante Abril 2026.",
             Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.95),
             font_size=13, color=C_GRIS_CLAR)

    footer_bar(sld)
    return sld


def slide_13_mapa_arquitectura(prs):
    """SLIDE 13 — Mapa visual de la arquitectura"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "La Nueva Arquitectura — Sistema de Dos Vías UTMACH",
               "2.919 artículos analizados → 8 unidades de investigación")

    # Título central
    add_rect(sld, Inches(2.5), Inches(1.3), Inches(8.3), Inches(0.5), C_AZUL_MED)
    add_text(sld, "SISTEMA DE INVESTIGACIÓN UTMACH · 2.919 artículos",
             Inches(2.6), Inches(1.35), Inches(8.1), Inches(0.42),
             font_size=14, bold=True, color=C_BLANCO, align=PP_ALIGN.CENTER)

    # VÍA A
    col_w = Inches(5.8)
    y0 = Inches(2.0)

    add_rect(sld, Inches(0.4), y0, col_w, Inches(0.5), C_AZUL_ACE)
    add_text(sld, "🔬  VÍA A — EXPERIMENTAL  (47% · 1.371 art.)",
             Inches(0.5), y0 + Inches(0.06), col_w - Inches(0.2), Inches(0.4),
             font_size=12, bold=True, color=C_BLANCO)

    centros_mapa = [
        "▸  Centro de Ciencias Químicas y Ambientales    [278 art. · 138.2 imp.]",
        "▸  Centro Experimental General                  [284 art. · 137.2 imp.]",
        "▸  Centro de Salud Integral y Biociencias       [265 art. · 126.6 imp.]",
        "▸  Centro Agroalimentario y Sostenibilidad      [178 art. ·  91.2 imp.]",
    ]
    y = y0 + Inches(0.55)
    for c in centros_mapa:
        add_rect(sld, Inches(0.4), y, col_w, Inches(0.56), C_AZUL_MED)
        add_text(sld, c, Inches(0.5), y + Inches(0.07), col_w - Inches(0.2), Inches(0.44),
                 font_size=11, color=C_GRIS_CLAR)
        y += Inches(0.6)

    add_rect(sld, Inches(0.4), y, col_w, Inches(0.5), C_AZUL_MED)
    add_text(sld, "💰  Alto presupuesto  ·  🔬  Laboratorios y equipos",
             Inches(0.5), y + Inches(0.07), col_w - Inches(0.2), Inches(0.38),
             font_size=11, color=C_AZUL_ACE, bold=True)

    # VÍA B
    x2 = Inches(7.1)
    add_rect(sld, x2, y0, col_w, Inches(0.5), C_VERDE_ACE)
    add_text(sld, "📊  VÍA B — SOCIAL/HUMANÍSTICA  (41.4% · 1.209 art.)",
             x2 + Inches(0.1), y0 + Inches(0.06), col_w - Inches(0.2), Inches(0.4),
             font_size=12, bold=True, color=C_BLANCO)

    obs_mapa = [
        "▸  Observatorio de Economía, Empresa e Innovación   [547 art. · 269.6 imp.]",
        "▸  Observatorio de Educación y Formación Profesional[241 art. · 114.4 imp.]",
        "▸  Observatorio de Derecho y Justicia Social        [127 art. ·  64.8 imp.]",
        "▸  Observatorio de Desarrollo Social y PP           [ 94 art. ·  50.5 imp.]",
    ]
    y = y0 + Inches(0.55)
    for o in obs_mapa:
        add_rect(sld, x2, y, col_w, Inches(0.56), C_AZUL_MED)
        add_text(sld, o, x2 + Inches(0.1), y + Inches(0.07), col_w - Inches(0.2), Inches(0.44),
                 font_size=11, color=C_GRIS_CLAR)
        y += Inches(0.6)

    add_rect(sld, x2, y, col_w, Inches(0.5), C_AZUL_MED)
    add_text(sld, "💡  Gestión ágil  ·  📊  Software, bases de datos, encuestas",
             x2 + Inches(0.1), y + Inches(0.07), col_w - Inches(0.2), Inches(0.38),
             font_size=11, color=C_VERDE_ACE, bold=True)

    # Separador VS
    add_rect(sld, Inches(6.55), y0, Inches(0.06), Inches(3.22), C_DORADO)

    footer_bar(sld)
    return sld


def slide_14_conclusion(prs):
    """SLIDE 14 — Conclusión & Call to Action"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)

    # Franja izquierda
    add_rect(sld, 0, 0, Inches(0.35), SLIDE_H, C_VERDE_ACE)

    add_text(sld, "Conclusión & Call to Action",
             Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.65),
             font_size=30, bold=True, color=C_BLANCO)
    add_text(sld, "La UTMACH lista para elevar su visibilidad internacional "
                  "financiando inteligentemente sus verdaderas fortalezas",
             Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.55),
             font_size=15, color=C_VERDE_ACE)
    add_rect(sld, Inches(0.6), Inches(1.5), Inches(10.0), Inches(0.04), C_DORADO)

    mensajes = [
        ("1", "Los datos hablan claro",
         "47% experimental + 41% social = dos ciencias distintas que merecen dos "
         "estrategias de inversión distintas. Sin ambigüedades.", C_AZUL_ACE),
        ("2", "La masa crítica existe",
         "Los investigadores élite y consolidados identificados por IPPC son suficientes "
         "para activar los 4 Centros y los 4 Observatorios sin reclutar externamente.", C_VERDE_ACE),
        ("3", "El sistema ya está diseñado",
         "Fases A y B: 100% completadas. La decisión que se toma hoy activa las Fases C y D "
         "hacia la reglamentación por el Consejo Universitario.", C_DORADO),
    ]
    y = Inches(1.7)
    for num, title, desc, clr in mensajes:
        add_rect(sld, Inches(0.6), y, Inches(0.65), Inches(0.95), clr)
        add_text(sld, num, Inches(0.6), y + Inches(0.18), Inches(0.65), Inches(0.55),
                 font_size=28, bold=True, color=C_BLANCO, align=PP_ALIGN.CENTER)
        add_rect(sld, Inches(1.3), y, Inches(11.5), Inches(0.95), C_AZUL_MED)
        add_text(sld, title, Inches(1.45), y + Inches(0.05), Inches(11.2), Inches(0.38),
                 font_size=14, bold=True, color=clr)
        add_text(sld, desc, Inches(1.45), y + Inches(0.45), Inches(11.2), Inches(0.44),
                 font_size=11, color=C_GRIS_CLAR, italic=True)
        y += Inches(1.05)

    # CTA
    add_rect(sld, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.15), C_AZUL_ACE)
    add_text(sld, "APROBEMOS HOY LA ARQUITECTURA DE DOS VÍAS",
             Inches(1.7), Inches(5.1), Inches(9.9), Inches(0.45),
             font_size=18, bold=True, color=C_BLANCO, align=PP_ALIGN.CENTER)
    add_text(sld,
             "Cruzar IPPC élite × clusters temáticos  →  Propuestas técnicas a Facultades  →  Reglamentación",
             Inches(1.7), Inches(5.58), Inches(9.9), Inches(0.45),
             font_size=13, color=C_DORADO, align=PP_ALIGN.CENTER)

    # Firma
    add_text(sld,
             '"No diseñamos centros. Descubrimos los que la UTMACH ya construyó con su producción científica."',
             Inches(0.6), Inches(6.38), Inches(12.1), Inches(0.55),
             font_size=13, italic=True, color=C_GRIS_TEXT, align=PP_ALIGN.CENTER)

    footer_bar(sld)
    return sld


def slide_15_backup(prs):
    """SLIDE 15 — Respaldo Técnico (Apéndice)"""
    sld = blank_slide(prs)
    fill_bg(sld, C_AZUL_OSC)
    header_bar(sld, "Respaldo Técnico — Apéndice Metodológico",
               "Para preguntas técnicas y auditoría del proceso")

    bloques = [
        ("CORPUS", [
            "2.919 artículos de revista (excl. libros, capítulos de libro)",
            "Variables: título · abstract · línea de investigación · cuartil Scopus/WoS · año",
            "87.3% de artículos sin cuartil Scopus/WoS — ponderación crítica",
        ], C_AZUL_ACE),
        ("PIPELINE NLP", [
            "TF-IDF con bigramas + stopwords académicas en español",
            "K-Means con Silhouette Score → K=20 óptimo",
            "Ponderación: Q1=1.0 · Q2=0.9 · Q3=0.8 · Q4=0.7 · Sin cuartil=0.2",
        ], C_VERDE_ACE),
        ("CLASIFICADOR EXPERIMENTAL", [
            "Léxico experimental: laboratorio, espectroscopía, cultivo, muestra...",
            "Léxico social: encuesta, política pública, análisis documental...",
            "Clasificación léxico-semántica sobre 2.919 artículos",
        ], C_DORADO),
        ("REPRODUCIBILIDAD", [
            "Scripts Python — versionados en GitHub (DIDI-UTMACH)",
            "master_analysis.py  ·  generate_figures.py  ·  generate_pptx.py",
            "Datos: produccion_clasificada_definitiva.csv · impacto_por_centro_definitivo.csv",
        ], C_ROJO_ACE),
    ]

    col_w = Inches(6.0)
    positions = [(Inches(0.4), Inches(1.35)), (Inches(6.9), Inches(1.35)),
                 (Inches(0.4), Inches(4.1)),  (Inches(6.9), Inches(4.1))]

    for (x, y), (titulo, items, clr) in zip(positions, bloques):
        add_rect(sld, x, y, col_w, Inches(0.38), clr)
        add_text(sld, titulo, x + Inches(0.1), y + Inches(0.05),
                 col_w - Inches(0.2), Inches(0.3), font_size=12, bold=True, color=C_BLANCO)
        yi = y + Inches(0.42)
        for item in items:
            add_rect(sld, x, yi, col_w, Inches(0.56), C_AZUL_MED)
            add_text(sld, "·  " + item, x + Inches(0.1), yi + Inches(0.07),
                     col_w - Inches(0.2), Inches(0.45), font_size=11, color=C_GRIS_CLAR)
            yi += Inches(0.6)

    footer_bar(sld)
    return sld


# ─── MAIN ────────────────────────────────────────────────────────────────────

def main():
    import sys
    # Forzar UTF-8 en stdout para evitar errores en terminales Windows cp1252
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    print("=" * 60)
    print("  Generando Presentacion_DIDI_UTMACH.pptx")
    print("=" * 60)

    prs = new_prs()

    builders = [
        slide_01_titulo,
        slide_02_problema,
        slide_03_metodologia,
        slide_04_epistemologia,
        slide_05_paradigma,
        slide_06_via_a,
        slide_07_via_b,
        slide_08_cuartiles,
        slide_09_ippc,
        slide_10_lineas,
        slide_11_dashboard,
        slide_12_roadmap,
        slide_13_mapa_arquitectura,
        slide_14_conclusion,
        slide_15_backup,
    ]

    for i, fn in enumerate(builders, 1):
        print(f"  Slide {i:02d}/{len(builders)}  -  {fn.__name__}")
        fn(prs)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print()
    print(f"[OK] Presentacion guardada en:")
    print(f"    {OUTPUT}")
    print()

    # Verificar figuras
    print("  Figuras insertadas:")
    for slide_n, path in sorted(FIGS.items()):
        status = "[OK]" if os.path.exists(path) else "[NO ENCONTRADA]"
        print(f"    Slide {slide_n:02d}: {os.path.basename(path)}  {status}")


if __name__ == "__main__":
    main()
